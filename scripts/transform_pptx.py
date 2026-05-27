#!/usr/bin/env python3
"""
transform_pptx.py
-----------------
Transforma una presentación docente antigua al formato T3 de UST.

Uso:
    python3 transform_pptx.py <archivo_antiguo.pptx> <plantilla_T3.pptx> <salida.pptx>

El script extrae el contenido del archivo antiguo y lo coloca en los layouts
correctos del template institucional, listo para que DG haga el diseño final.
"""

import sys
import copy
from dataclasses import dataclass, field
from typing import Optional
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from lxml import etree


# ── Tipos de diapositiva ──────────────────────────────────────────────────────
PORTADA      = "portada"       # L01 — Portada
INDICE       = "indice"        # L02 — Índice de contenidos
SECCION      = "seccion"       # L03 — Portada de sección
ICONOS_4     = "iconos_4"      # L04 — 4 ítems con ícono
LISTA_3      = "lista_3"       # L05 — 3 ítems con descripción
LISTA_IMG_3  = "lista_img_3"   # L06 — 3-4 ítems con imagen lateral
BLOQUES_4    = "bloques_4"     # L07 — 4 bloques con imagen
LISTA_IMG_4  = "lista_img_4"   # L08 — 4 ítems con imagen
LISTA_ALT    = "lista_alt"     # L09 — variante lista
TABLA        = "tabla"         # L10 — Tabla
PASOS_4      = "pasos_4"       # L11 — 4 pasos lineales
CIRCULOS_4   = "circulos_4"    # L12 — 4 círculos numerados
CIRCULOS_5   = "circulos_5"    # L13 — 5 círculos numerados
TIMELINE     = "timeline"      # L14 — Línea de tiempo
GRAFICO_1    = "grafico_1"     # L15 — Gráfico
GRAFICO_2    = "grafico_2"     # L16 — Gráfico variante
CONCLUSION   = "conclusion"    # L17 — Conclusión
REFERENCIAS  = "referencias"   # L18 — Referencias
CIERRE       = "cierre"        # L19 — Muchas gracias

# Mapeo tipo → índice en el template (0-based)
TEMPLATE_SLIDE_INDEX = {
    PORTADA:     0,
    INDICE:      1,
    SECCION:     2,
    ICONOS_4:    3,
    LISTA_3:     4,
    LISTA_IMG_3: 5,
    BLOQUES_4:   6,
    LISTA_IMG_4: 7,
    LISTA_ALT:   8,
    TABLA:       9,
    PASOS_4:     10,
    CIRCULOS_4:  11,
    CIRCULOS_5:  12,
    TIMELINE:    13,
    GRAFICO_1:   14,
    GRAFICO_2:   15,
    CONCLUSION:  16,
    REFERENCIAS: 17,
    CIERRE:      18,
}


@dataclass
class SlideContent:
    tipo: str
    titulo: str = ""
    subtitulo: str = ""
    items: list = field(default_factory=list)      # list of str
    descripcion: str = ""
    nota: str = ""


# ── Extracción de contenido ──────────────────────────────────────────────────

def extraer_textos_slide(slide) -> list[str]:
    """Retorna todos los textos no vacíos de una diapositiva."""
    textos = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                textos.append(t)
    return textos


def clasificar_slide(idx: int, textos: list[str], total: int) -> SlideContent:
    """Clasifica una diapositiva por su posición y contenido."""
    joined = " ".join(textos).lower()

    if idx == 0:
        # Portada
        titulo = textos[0] if textos else ""
        autor  = textos[1] if len(textos) > 1 else ""
        facult = textos[2] if len(textos) > 2 else ""
        return SlideContent(PORTADA, titulo=titulo, subtitulo=autor, descripcion=facult)

    if idx == total - 1 or "muchas gracias" in joined:
        return SlideContent(CIERRE)

    # Slides que son SOLO una URL o un link suelto → marcar como vacías para omitir
    if len(textos) == 1 and textos[0].startswith("http"):
        return SlideContent(SECCION, titulo="")  # Se omitirán

    if any(k in joined for k in ["referencia", "bibliograf"]):
        items = [t for t in textos if t]
        return SlideContent(REFERENCIAS, titulo="Referencias", items=items)

    if any(k in joined for k in ["contenidos", "contenido de la semana", "temas de"]) and idx < 3:
        titulo = textos[0] if textos else "Contenidos"
        items  = [t for t in textos[1:] if t]
        return SlideContent(INDICE, titulo=titulo, items=items)

    if any(k in joined for k in ["conclusión", "cierre", "ideas clave", "síntesis"]):
        titulo = textos[0] if textos else "Conclusión"
        items  = textos[1:]
        return SlideContent(CONCLUSION, titulo=titulo, items=items)

    # Slide vacía → omitir (se marca como SECCION vacía para filtrar)
    if not textos:
        return SlideContent(SECCION, titulo="")

    # Detectar portadas de sección: un solo texto grande o dos textos cortos
    if len(textos) <= 2 and all(len(t) < 60 for t in textos):
        return SlideContent(SECCION, titulo=textos[0], subtitulo=textos[1] if len(textos) > 1 else "")

    # Detectar slide de objetivo / aprendizaje
    if any(k in joined for k in ["objetivo", "aprendizaje", "pregunta de reflexión"]):
        return SlideContent(LISTA_3, titulo=textos[0], items=textos[1:])

    # Detectar slides con 4-5 ítems numerados (factores)
    items = [t for t in textos[1:] if t and len(t) < 120]
    if len(items) >= 4:
        titulo = textos[0] if textos else ""
        if len(items) >= 5:
            return SlideContent(CIRCULOS_5, titulo=titulo, items=items[:5])
        return SlideContent(CIRCULOS_4, titulo=titulo, items=items[:4])

    # Slide con 2-3 ítems descriptivos
    if 1 <= len(items) <= 3:
        titulo = textos[0] if textos else ""
        return SlideContent(LISTA_3, titulo=titulo, items=items)

    # Genérico: lista con descripción
    titulo = textos[0] if textos else ""
    return SlideContent(LISTA_IMG_3, titulo=titulo, items=textos[1:4])


def extraer_contenido(path_antiguo: str) -> list[SlideContent]:
    """Lee el PPTX antiguo y retorna la lista de SlideContent clasificados."""
    prs = Presentation(path_antiguo)
    total = len(prs.slides)
    slides_content = []

    for idx, slide in enumerate(prs.slides):
        textos = extraer_textos_slide(slide)
        sc = clasificar_slide(idx, textos, total)
        slides_content.append(sc)

    return slides_content


# ── Manipulación del template ─────────────────────────────────────────────────

def _replace_text_in_shape(shape, old: str, new: str):
    """Reemplaza texto en una shape preservando el primer formato encontrado."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
        # También reemplaza texto en el párrafo completo si no hay runs
        full = para.text
        if old in full and not para.runs:
            for elem in para._p:
                if elem.tag.endswith('}r'):
                    for t_elem in elem:
                        if t_elem.tag.endswith('}t') and old in t_elem.text:
                            t_elem.text = t_elem.text.replace(old, new)


def _set_shape_text(shape, text: str):
    """Establece el texto de una shape limpiando todo y poniendo el nuevo texto."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Usar el primer párrafo
    if tf.paragraphs:
        para = tf.paragraphs[0]
        if para.runs:
            para.runs[0].text = text
            # Limpiar runs adicionales
            for run in para.runs[1:]:
                run.text = ""
        else:
            # Crear un run si no hay
            from pptx.oxml.ns import qn
            r_elem = etree.SubElement(para._p, qn('a:r'))
            t_elem = etree.SubElement(r_elem, qn('a:t'))
            t_elem.text = text
    # Limpiar párrafos adicionales
    for para in tf.paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def copy_slide_from_template(template_prs: Presentation, slide_idx: int) -> object:
    """
    Duplica una diapositiva dentro del template_prs y la retorna.
    Trabaja siempre dentro del mismo objeto Presentation para conservar
    las referencias a imágenes y otros recursos embebidos.
    """
    source = template_prs.slides[slide_idx]
    layout = source.slide_layout

    # Crear nueva slide con el mismo layout
    new_slide = template_prs.slides.add_slide(layout)

    # Eliminar shapes por defecto que add_slide agrega
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copiar todos los elementos del spTree de origen (excepto los dos primeros: nvGrpSpPr y grpSpPr)
    src_tree = source.shapes._spTree
    new_tree = new_slide.shapes._spTree
    for elem in list(src_tree)[2:]:
        new_tree.append(copy.deepcopy(elem))

    return new_slide


PLACEHOLDERS_TITULO = ["Tema", "Título", "Título del recurso", "Conclusión", "Referencias"]
PLACEHOLDERS_ITEMS  = [
    "Lorem ipsum dolor sit amet, consectetur adipiscing elit. Suspendisse in mi sed velit lacinia vulputate.",
    "fsdfsfsf",
    "Sdgdsggd\nfsgsdgdgdg",
    "Sfsfjhskdf sfjshfshfhsjkdf\nsjkfhsjkhdgfshfd",
]
PLACEHOLDER_ITEM_PREFIX = "Ítem:"


def fill_portada(slide, sc: SlideContent):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if "Título del recurso" in t:
            _set_shape_text(shape, sc.titulo)
        elif "Nombre docente" in t:
            _set_shape_text(shape, sc.subtitulo)
        elif "Facultad:" in t:
            _set_shape_text(shape, sc.descripcion)


def fill_indice(slide, sc: SlideContent):
    items_restantes = list(sc.items)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if "Contenidos" in t:
            _set_shape_text(shape, sc.titulo or "Contenidos de la semana")
        elif t in ["Tema 1", "Tema 2", "Tema 3"]:
            if items_restantes:
                _set_shape_text(shape, items_restantes.pop(0))
            else:
                _set_shape_text(shape, "")
        elif any(ph in t for ph in ["Sdgdsggd", "Sfsfjhskdf", "sjkfhsjkhdgfshfd"]):
            _set_shape_text(shape, "")


def fill_seccion(slide, sc: SlideContent):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if "Título" in t and len(t) < 20:
            _set_shape_text(shape, sc.subtitulo or "")
        elif "Portada Tema" in t:
            _set_shape_text(shape, sc.titulo)


def fill_lista_items(slide, sc: SlideContent, max_items: int = 3):
    """Rellena slides tipo L05/L06 con título + ítems."""
    items = list(sc.items[:max_items])
    item_idx = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()

        if t in PLACEHOLDERS_TITULO or t == "Tema":
            _set_shape_text(shape, sc.titulo)
        elif t.startswith(PLACEHOLDER_ITEM_PREFIX) or t in ["fsdfsfsf", "Sdgdsggd\nfsgsdgdgdg"]:
            if item_idx < len(items):
                new_text = items[item_idx]
                item_idx += 1
                # Limpiar el prefijo "Ítem:" si el contenido ya lo tiene implícito
                _set_shape_text(shape, new_text)
            else:
                _set_shape_text(shape, "")


def fill_circulos(slide, sc: SlideContent, count: int = 4):
    """Rellena slides de círculos numerados (L12/L13)."""
    items = list(sc.items[:count])
    item_idx = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()

        if t in PLACEHOLDERS_TITULO or t in ["Tema", "Factores"]:
            _set_shape_text(shape, sc.titulo)
        elif t in [str(i) for i in range(1, count + 1)]:
            pass  # Mantener los números del círculo
        elif any(ph in t for ph in ["Lorem ipsum", "fsdfsfsf", "Sdgdsggd"]):
            if item_idx < len(items):
                _set_shape_text(shape, items[item_idx])
                item_idx += 1
            else:
                _set_shape_text(shape, "")


def fill_conclusion(slide, sc: SlideContent):
    items = list(sc.items)
    item_idx = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()

        if "Conclusión" in t:
            _set_shape_text(shape, sc.titulo or "Conclusión")
        elif t.startswith(PLACEHOLDER_ITEM_PREFIX):
            if item_idx < len(items):
                _set_shape_text(shape, items[item_idx])
                item_idx += 1
            else:
                _set_shape_text(shape, "")
        elif t in ["fsdfsfsf"]:
            _set_shape_text(shape, "")


def fill_referencias(slide, sc: SlideContent):
    items = list(sc.items)
    item_idx = 0

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()

        if "Referencias" in t:
            _set_shape_text(shape, "Referencias")
        elif t.startswith(PLACEHOLDER_ITEM_PREFIX):
            if item_idx < len(items):
                _set_shape_text(shape, items[item_idx])
                item_idx += 1
            else:
                _set_shape_text(shape, "")


def fill_slide(slide, sc: SlideContent):
    """Rellena una diapositiva del template con el contenido de SlideContent."""
    if sc.tipo == PORTADA:
        fill_portada(slide, sc)
    elif sc.tipo == INDICE:
        fill_indice(slide, sc)
    elif sc.tipo == SECCION:
        fill_seccion(slide, sc)
    elif sc.tipo in (LISTA_3, LISTA_IMG_3, LISTA_IMG_4, LISTA_ALT, ICONOS_4):
        fill_lista_items(slide, sc, max_items=4)
    elif sc.tipo in (CIRCULOS_4, PASOS_4):
        fill_circulos(slide, sc, count=4)
    elif sc.tipo == CIRCULOS_5:
        fill_circulos(slide, sc, count=5)
    elif sc.tipo == CONCLUSION:
        fill_conclusion(slide, sc)
    elif sc.tipo == REFERENCIAS:
        fill_referencias(slide, sc)
    # CIERRE no necesita modificación


# ── Script principal ──────────────────────────────────────────────────────────

def transformar(path_antiguo: str, path_template: str, path_salida: str):
    print(f"Leyendo archivo antiguo: {path_antiguo}")
    slides_content = extraer_contenido(path_antiguo)

    print(f"Diapositivas encontradas: {len(slides_content)}")
    for i, sc in enumerate(slides_content):
        print(f"  [{i+1:02d}] {sc.tipo:15s} | {sc.titulo[:50]}")

    print(f"\nCargando template: {path_template}")
    prs = Presentation(path_template)

    # El template tiene 19 slides originales (índices 0-18).
    # Duplicaremos las que necesitamos (se agregan al final) y luego
    # eliminaremos las 19 originales.
    TEMPLATE_ORIGINAL_COUNT = 19

    # Para cada contenido clasificado, duplicar el slide correcto del template
    generated_slides = []
    skipped = 0

    for i, sc in enumerate(slides_content):
        if sc.tipo == SECCION and not sc.titulo:
            skipped += 1
            continue

        template_idx = TEMPLATE_SLIDE_INDEX.get(sc.tipo, TEMPLATE_SLIDE_INDEX[LISTA_3])
        new_slide = copy_slide_from_template(prs, template_idx)
        fill_slide(new_slide, sc)
        generated_slides.append(new_slide)

    # Eliminar las 19 slides originales del template (están al inicio)
    xml_slides = prs.slides._sldIdLst
    for _ in range(TEMPLATE_ORIGINAL_COUNT):
        if len(xml_slides) > 0:
            xml_slides.remove(xml_slides[0])

    print(f"\nDiapositivas generadas: {len(prs.slides)} (omitidas vacías: {skipped})")
    prs.save(path_salida)
    print(f"Guardado: {path_salida}")


if __name__ == "__main__":
    if len(sys.argv) != 4:
        print("Uso: python3 transform_pptx.py <antiguo.pptx> <plantilla.pptx> <salida.pptx>")
        sys.exit(1)

    transformar(sys.argv[1], sys.argv[2], sys.argv[3])
